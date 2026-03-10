import sys
import os
import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

def extract_ppt_content(ppt_path, output_dir):
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"Error opening {ppt_path}: {e}")
        return None

    ppt_name = os.path.basename(ppt_path).replace('.pptx', '')
    assets_dir = os.path.join(output_dir, ppt_name, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    presentation_data = {
        "title": ppt_name,
        "slides": []
    }

    img_count = 0

    def process_shapes(shapes, slide, slide_idx, slide_data):
        nonlocal img_count
        for shape in shapes:
            if hasattr(shape, "shapes"):
                process_shapes(shape.shapes, slide, slide_idx, slide_data)
                continue

            if shape.has_text_frame:
                text = ""
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t: text += t + "\n"
                text = text.strip()
                if text:
                    if not slide_data["title"] and len(text) < 50:
                        slide_data["title"] = text
                    else:
                        slide_data["content"].append(text)
            
            try:
                xml = shape.element.xml
                matches = re.findall(r'r:embed="([^"]+)"', xml)
                for rId in matches:
                    try:
                        rel = slide.part.rels[rId]
                        image_part = rel.target_part
                        if "image" in image_part.content_type:
                            img_count += 1
                            blob = image_part.blob
                            
                            # Improved extension handling
                            content_type = image_part.content_type
                            if "png" in content_type: ext = "png"
                            elif "jpeg" in content_type: ext = "jpg"
                            elif "gif" in content_type: ext = "gif"
                            elif "svg" in content_type: ext = "svg"
                            elif "emf" in content_type: ext = "emf"
                            elif "wmf" in content_type: ext = "wmf"
                            else: ext = content_type.split('/')[-1].split('+')[0]
                            
                            img_filename = f"slide_{slide_idx+1}_img_{img_count}.{ext}"
                            img_path = os.path.join(assets_dir, img_filename)
                            with open(img_path, "wb") as f:
                                f.write(blob)
                            
                            web_path = f"/presentations/{ppt_name}/assets/{img_filename}"
                            if web_path not in slide_data["images"]:
                                slide_data["images"].append(web_path)
                    except KeyError:
                        continue 
            except Exception as e:
                print(f"Error processing shape XML on slide {slide_idx+1}: {e}")

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "index": slide_idx + 1,
            "title": "",
            "content": [],
            "images": []
        }
        process_shapes(slide.shapes, slide, slide_idx, slide_data)
        if not slide_data["title"]:
             slide_data["title"] = f"第 {slide_idx+1} 页"
        presentation_data["slides"].append(slide_data)
    
    json_path = os.path.join(output_dir, ppt_name, "data.json")
    with open(json_path, "w", encoding='utf-8') as f:
        json.dump(presentation_data, f, ensure_ascii=False, indent=4)
    
    return json_path

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python ppt_to_reveal.py <ppt_path> <output_dir>")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    out_dir = sys.argv[2]
    result = extract_ppt_content(ppt_file, out_dir)
    if result:
        print(f"Presentation data saved to {result}")
