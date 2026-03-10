from pptx import Presentation
from pptx.enum.dml import MSO_FILL
prs = Presentation("/Users/mdy/Desktop/宋/非遗课程/掐丝珐琅课件.pptx")
for i, slide in enumerate(prs.slides):
    if i >= 5: break
    print(f"Slide {i+1}:")
    for j, shape in enumerate(slide.shapes):
        try:
            fill_type = shape.fill.type if hasattr(shape, 'fill') else "N/A"
            print(f"  Shape {j+1}: type={shape.shape_type}, name={shape.name}, fill={fill_type}")
        except:
             print(f"  Shape {j+1}: Error getting info")
